from pathlib import Path
import os
from zipfile import BadZipFile
import xml.etree.ElementTree as ET

from pypdf import PdfReader

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    reader = PdfReader(path)
    parts: list[str] = []
    max_pages_raw = os.getenv("MAX_PDF_PAGES", "50").strip()
    try:
        max_pages = max(1, int(max_pages_raw))
    except ValueError:
        max_pages = 50

    for page in reader.pages[:max_pages]:
        extracted = page.extract_text()
        if extracted:
            parts.append(extracted)
    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PowerPoint uploads")

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                parts.append(text)
    return "\n".join(parts)


def _read_docx(path: Path) -> str:
    if Document is None:
        raise RuntimeError("python-docx is required for Word document uploads")

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    return "\n".join(parts)


def parse_document(file_path: str | Path) -> dict:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        text = _read_pdf(path)
    elif suffix == ".txt":
        text = _read_text_file(path)
    elif suffix == ".docx":
        try:
            text = _read_docx(path)
        except (BadZipFile, KeyError, RuntimeError) as exc:
            raise RuntimeError(f"Failed to read Word document: {exc}") from exc
    elif suffix == ".pptx":
        try:
            text = _read_pptx(path)
        except (BadZipFile, KeyError, RuntimeError) as exc:
            raise RuntimeError(f"Failed to read PowerPoint file: {exc}") from exc
    else:
        raise RuntimeError(f"Unsupported file type: {suffix or 'unknown'}")

    return {
        "file_id": path.stem.split("_", 1)[0] if "_" in path.stem else path.stem,
        "text": text,
        "page_count": None,
        "file_type": suffix.removeprefix("."),
    }
